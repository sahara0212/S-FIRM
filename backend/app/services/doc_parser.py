"""
업로드 문서에서 텍스트 추출
지원 형식: .pptx / .docx / .xlsx
"""
import io
from pathlib import Path


def extract_text(filename: str, content: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext == ".pptx":
        return _parse_pptx(content)
    elif ext == ".docx":
        return _parse_docx(content)
    elif ext in (".xlsx", ".xls"):
        return _parse_xlsx(content)
    else:
        # 평문 텍스트 시도
        try:
            return content.decode("utf-8")
        except Exception:
            return ""


def _parse_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            lines.append(f"[슬라이드 {i}]")
            lines.extend(slide_texts)
    return "\n".join(lines)


def _parse_docx(content: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(content))
    lines = []

    # 본문 단락
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # 스타일로 헤더 구분
            style = para.style.name if para.style else ""
            prefix = "## " if "Heading" in style else ""
            lines.append(prefix + text)

    # 테이블
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))

    return "\n".join(lines)


def _parse_xlsx(content: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"[시트: {sheet_name}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                lines.append(" | ".join(cells))
    return "\n".join(lines)
